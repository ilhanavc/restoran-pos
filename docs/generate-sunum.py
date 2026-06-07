# -*- coding: utf-8 -*-
"""
Restoran POS v1.1.0 — Kapsamlı Sunum (PPTX)
Her ekran görüntüsü için açıklamalı slayt + giriş/kapanış/mimari slaytlar.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

SCREENSHOTS = r"D:\dev\restoran-pos-v3\docs\screenshots"

# ── Renk Paleti ──────────────────────────────────────────────────────────────
INDIGO    = RGBColor(0x4F, 0x46, 0xE5)   # ana vurgu
INDIGO_LT = RGBColor(0xE0, 0xE7, 0xFF)   # açık zemin
DARK      = RGBColor(0x0F, 0x17, 0x2A)   # koyu başlık
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x6B, 0x72, 0x80)
GRAY_LT   = RGBColor(0xF3, 0xF4, 0xF6)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
AMBER     = RGBColor(0xD9, 0x77, 0x06)

# ── Slide boyutu: 16:9 Widescreen ────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # tamamen boş

# ── Yardımcı fonksiyonlar ─────────────────────────────────────────────────────
def rgb_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    rgb_fill(s, color)
    s.line.color.rgb = color if not line else RGBColor(0xD1, 0xD5, 0xDB)
    if not line:
        s.line.width = 0
    return s

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def add_image(slide, filename, l, t, w, h):
    path = os.path.join(SCREENSHOTS, filename)
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, l, t, width=w, height=h)
        # thin border
        pic.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        pic.line.width = Pt(0.5)
        return pic
    return None

def slide_number_badge(slide, n, total):
    """Sağ-alt köşe slayt numarası"""
    add_rect(slide, W - Inches(1.0), H - Inches(0.35),
             Inches(1.0), Inches(0.35), INDIGO)
    add_text(slide, f"{n} / {total}",
             W - Inches(1.0), H - Inches(0.35),
             Inches(1.0), Inches(0.35),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def bottom_bar(slide, label):
    """Alt şerit"""
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), DARK)
    add_text(slide, "Restoran POS v1.1.0  ·  İlhan AVCI  ·  220357008  ·  Atatürk Üniversitesi",
             Inches(0.2), H - Inches(0.35), W - Inches(1.2), Inches(0.35),
             size=9, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.LEFT)

# ── Veri: her slayt için içerik ───────────────────────────────────────────────
slides_data = [
    # KAPAK (özel)
    None,

    # 1: Giriş Ekranı
    {
        "title": "Giriş Ekranı",
        "subtitle": "Rol Bazlı Kimlik Doğrulama",
        "file": "01-login.png",
        "tag": "GÜVENLİK",
        "tag_color": INDIGO,
        "bullets": [
            "4 farklı kullanıcı rolü (Yönetici, Kasiyer, Garson, Mutfak)",
            "JWT tabanlı oturum yönetimi (1 saat access token)",
            "Demo hızlı giriş butonları ile tek tıkla erişim",
            "Şifre politikası: min 8 karakter + büyük harf + rakam",
            "Zorunlu şifre değiştirme akışı (admin reset sonrası)",
            "'Şifremi Unuttum' sıfırlama talep akışı",
        ],
        "highlight": "bcrypt cost 10 · JWT · Refresh Token · Rate Limit",
    },

    # 2: Masa Yönetimi
    {
        "title": "Masa Yönetimi",
        "subtitle": "Görsel Alan ve Doluluk Takibi",
        "file": "02-tables.png",
        "tag": "OPERASYON",
        "tag_color": GREEN,
        "bullets": [
            "4 bölge (İç Salon, Bahçe, VIP, Üst Kat), 24 masa",
            "Renk skalası ile anlık doluluk durumu",
            "Masalar arası sipariş transferi",
            "Masa kapasitesi ve alan yönetimi",
            "Socket.io ile anlık durum güncellemeleri",
            "Yoğun saatlerde hızlı masa seçimi optimizasyonu",
        ],
        "highlight": "Gerçek zamanlı senkronizasyon — Socket.io",
    },

    # 3: Sipariş Alma
    {
        "title": "Sipariş Akışı",
        "subtitle": "Hızlı Kategori → Ürün → Sepet Süreci",
        "file": "03-order.png",
        "tag": "OPERASYON",
        "tag_color": GREEN,
        "bullets": [
            "Sol: Kategori listesi | Orta: Ürünler | Sağ: Sepet",
            "Her ürün için modifier desteği (acılı, az pişmiş vb.)",
            "Kalem bazlı özel not girişi",
            "'Mutfağa Gönder' → Socket.io ile anlık mutfak güncellemesi",
            "Sipariş düzenleme (adet, not, silme)",
            "Yazıcı routing: her kategori farklı yazıcıya yönlendirilebilir",
        ],
        "highlight": "Sipariş → Mutfak gecikmesi: p95 < 100 ms",
    },

    # 4: Ödeme
    {
        "title": "Ödeme Ekranı",
        "subtitle": "Çoklu Ödeme Tipi ve Hızlı İşlem",
        "file": "04-payment.png",
        "tag": "OPERASYON",
        "tag_color": GREEN,
        "bullets": [
            "Nakit / Kart / Karışık ödeme desteği",
            "Hızlı tutar butonları: 50 / 100 / 200 / 500 TL",
            "Para üstü otomatik hesaplaması",
            "İndirim uygulama (tutar veya yüzde)",
            "Ödeme tamamlanınca masa otomatik boşa düşer",
            "Kasa fişi ESC/POS PC857 ile otomatik basılır",
        ],
        "highlight": "Ödeme → fiş: tam atomic transaction",
    },

    # 5: Mutfak
    {
        "title": "Mutfak Ekranı",
        "subtitle": "Anlık Sipariş Takip ve Yaş Uyarıları",
        "file": "05-kitchen.png",
        "tag": "GERÇEK ZAMANLI",
        "tag_color": AMBER,
        "bullets": [
            "Aktif siparişler kart formatında listelenir",
            "Yaş renk kodları: ⚪ 0-10 dk → 🟡 10-20 dk → 🔴 20+ dk",
            "Kalem bazlı 'Hazırlandı' işaretleme",
            "Socket.io: yeni sipariş → ekran anında güncellenir (yenileme yok)",
            "Mutfak personeline özel sadeleştirilmiş arayüz",
            "Sesli bildirim (yeni sipariş geldiğinde)",
        ],
        "highlight": "Polling yok — saf Socket.io real-time",
    },

    # 6: Paket Sipariş
    {
        "title": "Paket Sipariş + CallerID",
        "subtitle": "Otomatik Müşteri Tanıma",
        "file": "06-takeaway.png",
        "tag": "DONANIM ENT.",
        "tag_color": RGBColor(0x0E, 0x90, 0xD2),
        "bullets": [
            "CallerID HID cihaz veya PowerShell pano dinleme",
            "Kayıtlı müşteri arıyorsa adı ve adresi otomatik açılır",
            "Yeni müşteri için hızlı kayıt formu",
            "Çoklu teslimat adresi seçimi",
            "Ödeme tipi seçimi zorunlu (nakit/kart — raporlar için)",
            "Mutfak fişi otomatik basılır",
        ],
        "highlight": "CallerID: C# .NET 8 self-contained binary",
    },

    # 7: Müşteriler
    {
        "title": "Müşteri Yönetimi",
        "subtitle": "360° Müşteri Profili",
        "file": "07-customers.png",
        "tag": "CRM",
        "tag_color": RGBColor(0x7C, 0x3A, 0xED),
        "bullets": [
            "Sayfalama: 50 kayıt/sayfa, 'Daha Fazla Yükle'",
            "Arama: ad, soyad, telefon (normalize edilmiş)",
            "360° profil: toplam harcama, sipariş sayısı, son ziyaret",
            "En çok sipariş ettiği 3 ürün",
            "Çoklu telefon ve adres yönetimi",
            "Excel / CSV toplu import-export",
        ],
        "highlight": "Telefon normalize: E.164 benzeri dedupe",
    },

    # 8: Raporlar
    {
        "title": "Raporlar ve Analizler",
        "subtitle": "4 İnteraktif Grafik + X/Z Dönem Kapatma",
        "file": "08-reports.png",
        "tag": "ANALİTİK",
        "tag_color": RGBColor(0xDB, 0x27, 0x77),
        "bullets": [
            "📈 Saatlik satış dağılımı (line chart)",
            "🥧 Kategori dağılımı (pie chart)",
            "🍩 Ödeme tipi oranı / nakit vs kart (donut)",
            "📊 En çok satan ürünler top-10 (bar chart)",
            "X Raporu (ara) / Z Raporu (gün sonu kapatma)",
            "Excel ve PDF dışa aktarma",
        ],
        "highlight": "Recharts + X/Z period close (kapalı döneme iade yasağı)",
    },

    # 9: Menü Tanımları
    {
        "title": "Menü Tanımları",
        "subtitle": "Esnek Kategori ve Ürün Yönetimi",
        "file": "09-menu-mgmt.png",
        "tag": "YÖNETİM",
        "tag_color": RGBColor(0x05, 0x96, 0x69),
        "bullets": [
            "Sol panel: kategoriler | Sağ panel: ürünler",
            "Kategori: ikon, renk, yazıcı hedefi seçimi",
            "Ürün: resim, modifier grupları, kombo menü",
            "Akıllı silme: geçmiş kayıt varsa soft-delete, yoksa hard-delete",
            "Kategori bazlı printer routing (mutfak / bar / kasa)",
            "Sıralama (drag-and-drop benzeri sort_order)",
        ],
        "highlight": "Hybrid delete: veri bütünlüğü korunur",
    },

    # 10: Ayarlar
    {
        "title": "Ayarlar Merkezi",
        "subtitle": "Tek Ekranda Tüm Yapılandırma",
        "file": "10-settings.png",
        "tag": "SİSTEM",
        "tag_color": GRAY,
        "bullets": [
            "İşletme bilgileri (ad, telefon, adres, vergi no)",
            "Kullanıcı yönetimi: ekle, rol değiştir, şifre sıfırla",
            "Yazıcı tanımları ve IP/USB yapılandırması",
            "Store Bridge durumu ve log görüntüleme",
            "Sürüm notları ve güncelleme kontrolü",
            "Cihaz eşleştirme (QR kod ile mobil pairing)",
        ],
        "highlight": "electron-updater ile otomatik güncelleme",
    },

    # 11: Yedekleme
    {
        "title": "Yedekleme ve Geri Yükleme",
        "subtitle": "SHA-256 Doğrulamalı Güvenli Yedek",
        "file": "11-backup.png",
        "tag": "VERİ GÜVENLİĞİ",
        "tag_color": RGBColor(0xDC, 0x26, 0x26),
        "bullets": [
            "Otomatik yedek: gece 02:00 (cron)",
            "İçerik: pos.db + uploads + pos-config.json",
            "meta.json: rowCounts, sha256, integrityCheck, appVersion",
            "30 gün otomatik rotasyon",
            "Geri yükleme: 2 adımlı modal + SHA-256 bütünlük kontrolü",
            "Başarısızsa otomatik safety revert",
        ],
        "highlight": "RPO: 24 saat · RTO: 15 dakika",
    },

    # 12: Audit Log
    {
        "title": "Denetim Kayıtları",
        "subtitle": "Kim, Ne, Ne Zaman — Tam İzlenebilirlik",
        "file": "12-audit-log.png",
        "tag": "DENETİM",
        "tag_color": DARK,
        "bullets": [
            "Tüm önemli değişiklikler entity_mutations tablosunda saklanır",
            "Before / After JSON karşılaştırma görünümü",
            "Filtreler: tablo, işlem tipi, tarih, kullanıcı",
            "Ürün, sipariş, ödeme, iade, müşteri, kullanıcı kayıtları",
            "Immutable append-only (kayıt silinemez, değiştirilemez)",
            "Mali denetim ve sorumluluk takibi için kritik",
        ],
        "highlight": "Audit trail: OWASP A09 — Security Logging uyumu",
    },

    # TEKNOLOJİ (özel)
    None,

    # METRİKLER (özel)
    None,

    # KAPANIŞ (özel)
    None,
]

TOTAL = len(slides_data)  # gerçek toplam aşağıda düzeltilecek


# ─────────────────────────────────────────────────────────────────────────────
# SLAYT 1 — KAPAK
# ─────────────────────────────────────────────────────────────────────────────
def make_cover(prs):
    sld = prs.slides.add_slide(BLANK)

    # Sol koyu blok
    add_rect(sld, 0, 0, Inches(5.2), H, DARK)

    # Sağ açık zemin
    add_rect(sld, Inches(5.2), 0, W - Inches(5.2), H, INDIGO_LT)

    # İkon arka planı (indigo daire efekti)
    add_rect(sld, Inches(0.4), Inches(1.2), Inches(0.7), Inches(0.7),
             RGBColor(0x6D, 0x64, 0xFF))

    # "P" harfi
    add_text(sld, "P", Inches(0.4), Inches(1.2), Inches(0.7), Inches(0.7),
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Ana başlık
    add_text(sld, "RESTORAN POS",
             Inches(0.3), Inches(2.2), Inches(4.7), Inches(1.2),
             size=40, bold=True, color=WHITE)

    add_text(sld, "Yazılım Gereksinimleri ve Özellik Sunumu",
             Inches(0.3), Inches(3.5), Inches(4.7), Inches(0.7),
             size=15, color=RGBColor(0xA5, 0xB4, 0xFC), italic=True)

    add_text(sld, "Sürüm 1.1.0  ·  Mayıs 2026",
             Inches(0.3), Inches(4.3), Inches(4.7), Inches(0.4),
             size=12, color=RGBColor(0x6B, 0x72, 0x80))

    # Ayraç
    add_rect(sld, Inches(0.3), Inches(5.0), Inches(2.0), Inches(0.04),
             RGBColor(0x4F, 0x46, 0xE5))

    # Yazar
    add_text(sld, "İlhan AVCI  |  220357008",
             Inches(0.3), Inches(5.2), Inches(4.7), Inches(0.4),
             size=13, bold=True, color=WHITE)

    add_text(sld, "Atatürk Üniversitesi\nİktisadi ve İdari Bilimler Fakültesi\nYönetim Bilişim Sistemleri Bölümü",
             Inches(0.3), Inches(5.7), Inches(4.7), Inches(1.2),
             size=11, color=RGBColor(0x9C, 0xA3, 0xAF))

    # Sağ taraf — ekran görüntüsü montajı (3 küçük ekran)
    shots = [("01-login.png", Inches(5.5), Inches(0.4)),
             ("03-order.png",  Inches(8.0), Inches(0.4)),
             ("08-reports.png",Inches(5.5), Inches(3.9)),
             ("02-tables.png", Inches(8.0), Inches(3.9))]
    for fname, lx, ty in shots:
        add_image(sld, fname, lx, ty, Inches(2.3), Inches(3.0))

    # Sağ alt dekoratif
    add_text(sld, "Bitirme Projesi",
             Inches(5.5), H - Inches(0.4), Inches(4.0), Inches(0.4),
             size=10, color=RGBColor(0x4F, 0x46, 0xE5),
             align=PP_ALIGN.CENTER, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLAYT: Ekran görüntülü içerik slaytı
# ─────────────────────────────────────────────────────────────────────────────
def make_content_slide(prs, data, slide_num, total):
    sld = prs.slides.add_slide(BLANK)

    # Arka plan
    add_rect(sld, 0, 0, W, H, GRAY_LT)

    # Sol panel (koyu)
    add_rect(sld, 0, 0, Inches(4.5), H, DARK)

    # Üst renk şeridi (tag rengi)
    add_rect(sld, 0, 0, Inches(4.5), Inches(0.08), data["tag_color"])

    # Tag (etiket)
    add_rect(sld, Inches(0.3), Inches(0.25),
             Inches(1.8), Inches(0.28), data["tag_color"])
    add_text(sld, data["tag"],
             Inches(0.3), Inches(0.25), Inches(1.8), Inches(0.28),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Başlık
    add_text(sld, data["title"],
             Inches(0.25), Inches(0.65), Inches(4.1), Inches(0.9),
             size=28, bold=True, color=WHITE)

    # Alt başlık
    add_text(sld, data["subtitle"],
             Inches(0.25), Inches(1.5), Inches(4.1), Inches(0.4),
             size=13, color=RGBColor(0xA5, 0xB4, 0xFC), italic=True)

    # Ayraç
    add_rect(sld, Inches(0.25), Inches(1.95), Inches(1.5), Inches(0.035),
             data["tag_color"])

    # Bullets
    bullet_top = Inches(2.15)
    for i, b in enumerate(data["bullets"]):
        # bullet ikon
        add_rect(sld, Inches(0.25), bullet_top + Inches(i * 0.58) + Inches(0.08),
                 Inches(0.07), Inches(0.07), data["tag_color"])
        add_text(sld, b,
                 Inches(0.4), bullet_top + Inches(i * 0.58),
                 Inches(3.95), Inches(0.55),
                 size=11, color=WHITE, wrap=True)

    # Highlight kutusu (alt)
    add_rect(sld, Inches(0.2), H - Inches(0.9), Inches(4.1), Inches(0.55),
             RGBColor(0x1E, 0x29, 0x3B))
    add_text(sld, "🔧  " + data["highlight"],
             Inches(0.3), H - Inches(0.88), Inches(4.0), Inches(0.5),
             size=10, color=RGBColor(0xA5, 0xB4, 0xFC), italic=True, wrap=True)

    # Sağ: ekran görüntüsü
    add_rect(sld, Inches(4.7), Inches(0.15), Inches(8.45), Inches(6.85),
             WHITE)
    add_image(sld, data["file"],
              Inches(4.85), Inches(0.3), Inches(8.1), Inches(6.55))

    # Slayt numarası
    slide_number_badge(sld, slide_num, total)

    return sld


# ─────────────────────────────────────────────────────────────────────────────
# TEKNOLOJİ slaytı
# ─────────────────────────────────────────────────────────────────────────────
def make_tech_slide(prs, slide_num, total):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, W, H, DARK)
    add_rect(sld, 0, 0, W, Inches(0.06), INDIGO)

    add_text(sld, "Teknoloji Yığını", Inches(0.5), Inches(0.25),
             Inches(8), Inches(0.7), size=32, bold=True, color=WHITE)
    add_text(sld, "Modern, test güdümlü ve güvenli bir yazılım mimarisi",
             Inches(0.5), Inches(0.95), Inches(10), Inches(0.4),
             size=13, color=RGBColor(0xA5, 0xB4, 0xFC), italic=True)

    cols = [
        ("Frontend", INDIGO,
         ["React 18 + Vite 5", "React Router 6 (HashRouter)", "Recharts", "Tailwind CSS"]),
        ("Backend", RGBColor(0x05, 0x96, 0x69),
         ["Node.js 20 + Express 4", "Socket.io 4", "Zod validation", "jsonwebtoken + bcryptjs"]),
        ("Veri & Depolama", AMBER,
         ["better-sqlite3 (WAL)", "13 migration (forward-only)", "Dual-write _cents", "Snapshot kolonları"]),
        ("Masaüstü & Paket", RGBColor(0xDC, 0x26, 0x26),
         ["Electron 34", "electron-builder 24", "electron-updater", "C# .NET 8 CallerID"]),
        ("Gözlemlenebilirlik", RGBColor(0x7C, 0x3A, 0xED),
         ["Pino NDJSON logging", "Sentry + redact", "X-Request-Id corr.", "Crash reporter"]),
        ("Test & Kalite", RGBColor(0x0E, 0x90, 0xD2),
         ["Vitest + Supertest", "Playwright E2E", "GitHub Actions CI", "ESLint 0 warning"]),
    ]

    col_w = Inches(2.0)
    gap   = Inches(0.12)
    start_l = Inches(0.4)
    for i, (label, color, items) in enumerate(cols):
        lx = start_l + i * (col_w + gap)
        # header box
        add_rect(sld, lx, Inches(1.55), col_w, Inches(0.38), color)
        add_text(sld, label, lx, Inches(1.55), col_w, Inches(0.38),
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # content box
        add_rect(sld, lx, Inches(1.93), col_w,
                 Inches(5.0), RGBColor(0x1C, 0x24, 0x3B))
        for j, item in enumerate(items):
            add_rect(sld, lx + Inches(0.1), Inches(2.05) + j * Inches(1.1),
                     Inches(0.08), Inches(0.08), color)
            add_text(sld, item,
                     lx + Inches(0.25), Inches(2.0) + j * Inches(1.1),
                     col_w - Inches(0.35), Inches(1.0),
                     size=10, color=WHITE, wrap=True)

    slide_number_badge(sld, slide_num, total)


# ─────────────────────────────────────────────────────────────────────────────
# METRİKLER slaytı
# ─────────────────────────────────────────────────────────────────────────────
def make_metrics_slide(prs, slide_num, total):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, W, H, GRAY_LT)
    add_rect(sld, 0, 0, W, Inches(1.2), DARK)

    add_text(sld, "Kalite Metrikleri ve Proje Kapsamı",
             Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
             size=30, bold=True, color=WHITE)
    add_text(sld, "452 otomatik test · 0 lint uyarı · 9.3/10 genel kalite skoru",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.35),
             size=13, color=RGBColor(0xA5, 0xB4, 0xFC), italic=True)

    # Büyük metrik kartlar
    big_metrics = [
        ("452", "Otomatik Test", INDIGO),
        ("0", "Lint Uyarı", GREEN),
        ("13", "Migration", AMBER),
        ("9.3/10", "Kalite Skoru", RGBColor(0xDC, 0x26, 0x26)),
        ("30+", "Tablo", RGBColor(0x7C, 0x3A, 0xED)),
        ("60+", "API Endpoint", RGBColor(0x0E, 0x90, 0xD2)),
    ]

    card_w = Inches(2.0)
    card_h = Inches(1.6)
    gap = Inches(0.18)
    start_l = Inches(0.4)
    top = Inches(1.4)

    for i, (val, label, color) in enumerate(big_metrics):
        lx = start_l + i * (card_w + gap)
        add_rect(sld, lx, top, card_w, card_h, WHITE)
        add_rect(sld, lx, top, card_w, Inches(0.06), color)
        add_text(sld, val, lx, top + Inches(0.2), card_w, Inches(0.9),
                 size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(sld, label, lx, top + Inches(1.05), card_w, Inches(0.45),
                 size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Alt kısım: Sprint listesi 2 kolon
    sprints_left = [
        "Sprint 1-3: Güvenlik, Yazıcı, UX Temelleri",
        "Sprint 4-6: Yeni Özellikler + Real-Time",
        "Sprint 7-8: Test Kapsamı + Prod Hardening",
        "Sprint 9-10: Feature Completion + Audit",
        "Sprint 11-12: Backup/Restore + Repo Temizlik",
    ]
    sprints_right = [
        "D-1...D-5: CI, Sertleştirme, Decomposition",
        "DB-1...DB-4: Migration, Audit Trail, Para, Snapshot",
        "C-1, C-2: Cloud Desteği + Refresh Token",
        "M-1.1, M-1.2: Mobil Endpoint + Cihaz Pairing",
        "FAZ 0 (0.1-0.9): Tüm Güvenlik Önlemleri",
    ]

    add_text(sld, "Geliştirme Yol Haritası",
             Inches(0.4), Inches(3.3), Inches(5), Inches(0.4),
             size=14, bold=True, color=DARK)

    for i, s in enumerate(sprints_left):
        add_rect(sld, Inches(0.4), Inches(3.8) + i * Inches(0.52),
                 Inches(0.1), Inches(0.1), INDIGO)
        add_text(sld, s, Inches(0.6), Inches(3.75) + i * Inches(0.52),
                 Inches(6.0), Inches(0.5), size=11, color=DARK, wrap=True)

    for i, s in enumerate(sprints_right):
        add_rect(sld, Inches(6.8), Inches(3.8) + i * Inches(0.52),
                 Inches(0.1), Inches(0.1), GREEN)
        add_text(sld, s, Inches(7.0), Inches(3.75) + i * Inches(0.52),
                 Inches(6.0), Inches(0.5), size=11, color=DARK, wrap=True)

    slide_number_badge(sld, slide_num, total)


# ─────────────────────────────────────────────────────────────────────────────
# KAPANIŞ slaytı
# ─────────────────────────────────────────────────────────────────────────────
def make_closing_slide(prs, slide_num, total):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, W, H, DARK)

    # Dekoratif çizgi
    add_rect(sld, Inches(3.0), Inches(2.3), Inches(7.0), Inches(0.05), INDIGO)

    add_text(sld, "Teşekkürler",
             0, Inches(1.2), W, Inches(1.2),
             size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sld, "İlhan AVCI  ·  220357008",
             0, Inches(2.55), W, Inches(0.6),
             size=18, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC),
             align=PP_ALIGN.CENTER)

    add_text(sld,
             "Atatürk Üniversitesi  ·  İktisadi ve İdari Bilimler Fakültesi\n"
             "Yönetim Bilişim Sistemleri Bölümü  ·  Bitirme Projesi  ·  Mayıs 2026",
             0, Inches(3.2), W, Inches(0.9),
             size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Erişim linkleri kartı
    add_rect(sld, Inches(3.5), Inches(4.2), Inches(6.2), Inches(2.0),
             RGBColor(0x1C, 0x24, 0x3B))
    add_rect(sld, Inches(3.5), Inches(4.2), Inches(6.2), Inches(0.06), INDIGO)

    links = [
        ("📦 Uygulama:", "github.com/ilhanavc/restoran-pos/releases/tag/v1.1.0-graduation"),
        ("🐙 Kaynak Kod:", "github.com/ilhanavc/restoran-pos"),
        ("📧 İletişim:", "ilhanavci499@gmail.com"),
    ]
    for i, (label, val) in enumerate(links):
        add_text(sld, label,
                 Inches(3.7), Inches(4.4) + i * Inches(0.55),
                 Inches(1.8), Inches(0.5),
                 size=11, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC))
        add_text(sld, val,
                 Inches(5.5), Inches(4.4) + i * Inches(0.55),
                 Inches(4.0), Inches(0.5),
                 size=11, color=WHITE)

    slide_number_badge(sld, slide_num, total)


# ─────────────────────────────────────────────────────────────────────────────
# ÇALIŞTIR
# ─────────────────────────────────────────────────────────────────────────────
content_count = sum(1 for d in slides_data if d is not None)
TOTAL_SLIDES = 1 + content_count + 3  # kapak + içerik + tech + metrics + closing

# Kapak
make_cover(prs)

# İçerik slaytları
slide_num = 2
for data in slides_data:
    if data is None:
        continue
    make_content_slide(prs, data, slide_num, TOTAL_SLIDES)
    slide_num += 1

# Özel slaytlar
make_tech_slide(prs, slide_num, TOTAL_SLIDES);     slide_num += 1
make_metrics_slide(prs, slide_num, TOTAL_SLIDES);  slide_num += 1
make_closing_slide(prs, slide_num, TOTAL_SLIDES)

out = r"D:\dev\restoran-pos-v3\docs\SUNUM-Restoran-POS-v1.1.0.pptx"
prs.save(out)
import os
print(f"OK — {TOTAL_SLIDES} slayt → {os.path.getsize(out)/1024/1024:.1f} MB")
print(f"Dosya: {out}")
